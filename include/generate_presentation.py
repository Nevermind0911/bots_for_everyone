#!/usr/bin/env python3
"""
Генерация PPTX-презентации для заявки на конкурс «Студенческий Стартап» (VII очередь).
Проект: BotWork — Telegram-боты для бизнеса.

Использование:
    pip install python-pptx
    python generate_presentation.py

Результат: файл BotWork_Application.pptx в текущей директории.
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ─── Пути ─────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).resolve().parent          # include/
PROJECT_DIR = SCRIPT_DIR.parent                        # python_project_site/
IMAGES_DIR = SCRIPT_DIR / "images"
OUTPUT_DIR = PROJECT_DIR / "docs"
IMG_LEFT = str(IMAGES_DIR / "picture_left_corner.png")
IMG_RIGHT_FIRST_LAST = str(IMAGES_DIR / "picture_right_corner_first_last.png")
IMG_RIGHT_TEXT = str(IMAGES_DIR / "picture_right_corner_text_slides.png")

# ─── Константы стиля (из шаблона Application_task_description.pptx) ───
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

# Палитра шаблона
BLUE = RGBColor(0xB5, 0xD8, 0xE1)
PINK = RGBColor(0xFB, 0xB3, 0xC1)
RED = RGBColor(0xFD, 0x49, 0x3D)
ORANGE = RGBColor(0xFC, 0xAF, 0x17)
PURPLE = RGBColor(0x8F, 0x00, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
GRAY = RGBColor(0x33, 0x33, 0x33)

FONT_MAIN = "Georgia"
FONT_ALT = "Arial"


# ─── Утилиты ─────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, text="",
             font_name=FONT_MAIN, font_size=Pt(24), font_bold=True,
             font_color=WHITE, alignment=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = font_size
        p.font.bold = font_bold
        p.font.color.rgb = font_color
        p.alignment = alignment
    return shape


def add_textbox(slide, left, top, width, height, paragraphs,
                font_name=FONT_MAIN, font_size=Pt(16), font_color=DARK,
                line_spacing=1.3, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, bold, size_override, color_override) in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = font_name
        p.font.size = size_override or font_size
        p.font.bold = bold
        p.font.color.rgb = color_override or font_color
        p.alignment = alignment
        p.space_after = Pt(4)

    return txBox


def add_color_bar(slide, top, colors=None):
    """Декоративная цветная полоса под заголовком — как в шаблоне."""
    if colors is None:
        colors = [PINK, PURPLE, RED, ORANGE, BLUE]
    bar_h = Emu(100000)
    seg_w = Emu(1800000)
    for i, c in enumerate(colors):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(i * int(seg_w)), top, seg_w, bar_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()


def add_section_header(slide, section_name, header_color=RED,
                       bar_top=Emu(520000)):
    """Цветной прямоугольник-заголовок секции + декоративная полоса."""
    add_rect(slide, Emu(0), Emu(30000), Emu(5500000), Emu(470000),
             header_color, section_name,
             font_size=Pt(24), font_bold=True, font_color=WHITE)
    add_color_bar(slide, bar_top)


def add_corner_images(slide, first_last=False):
    """Логотип ФСИ (лев. низ) + картинка (прав. низ).

    first_last=True  → picture_right_corner_first_last.png  (титульный/финальный)
    first_last=False → picture_right_corner_text_slides.png (контентные слайды)
    """
    if os.path.isfile(IMG_LEFT):
        slide.shapes.add_picture(
            IMG_LEFT,
            Emu(256000), Emu(5770000),   # нижний левый угол
            Emu(1670000), Emu(820000),
        )
    img_right = IMG_RIGHT_FIRST_LAST if first_last else IMG_RIGHT_TEXT
    if os.path.isfile(img_right):
        slide.shapes.add_picture(
            img_right,
            Emu(7380000), Emu(5370000),  # нижний правый угол
            Emu(1634000), Emu(1288000),
        )


def content_slide(prs, section_name, header_color, paragraphs):
    """Создаёт стандартный слайд с заголовком секции и текстовым блоком."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)
    add_section_header(slide, section_name, header_color)
    add_textbox(slide, Emu(460000), Emu(750000),
                Emu(8200000), Emu(5500000), paragraphs)
    add_corner_images(slide)
    return slide


# ─── Генерация презентации ────────────────────────────────────────────
def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 1 — Титульный
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Декоративные блоки
    add_rect(slide, Emu(0), Emu(3900000), Emu(5700000), Emu(2950000), BLUE)
    add_rect(slide, Emu(5700000), Emu(5900000), Emu(3450000), Emu(960000), PINK)
    add_rect(slide, Emu(5700000), Emu(3900000), Emu(800000), Emu(2000000), RED)
    add_rect(slide, Emu(6500000), Emu(3900000), Emu(2650000), Emu(2000000), PURPLE)
    add_rect(slide, Emu(6500000), Emu(0), Emu(2650000), Emu(3900000), ORANGE)

    # Заголовок
    add_textbox(slide, Emu(800000), Emu(1200000), Emu(5000000), Emu(2500000),
                [
                    ("BotWork", True, Pt(44), PURPLE),
                    ("", False, Pt(8), None),
                    ("Платформа для создания", False, Pt(24), DARK),
                    ("Telegram-ботов для бизнеса", False, Pt(24), DARK),
                ],
                alignment=PP_ALIGN.LEFT)

    # Метаданные
    add_textbox(slide, Emu(300000), Emu(4100000), Emu(5000000), Emu(2500000),
                [
                    ("Номер заявки: ___________", False, Pt(16), WHITE),
                    ("", False, Pt(6), None),
                    ("Заявитель: Репин Артемий, Мартынов Борис, Харисов Денис", False, Pt(16), WHITE),
                ],
                alignment=PP_ALIGN.LEFT)

    # «Студенческий Стартап»
    add_rect(slide, Emu(6700000), Emu(1500000), Emu(2300000), Emu(600000),
             PURPLE, "Студенческий Стартап",
             font_name=FONT_ALT, font_size=Pt(16),
             font_bold=True, font_color=WHITE,
             alignment=PP_ALIGN.CENTER)

    add_rect(slide, Emu(7000000), Emu(2200000), Emu(1700000), Emu(420000),
             RED, "VII очередь",
             font_name=FONT_ALT, font_size=Pt(20),
             font_bold=False, font_color=WHITE,
             alignment=PP_ALIGN.CENTER)

    # Логотип ФСИ — верхний левый угол титульного слайда
    if os.path.isfile(IMG_LEFT):
        slide.shapes.add_picture(
            IMG_LEFT,
            Emu(266000), Emu(333000),
            Emu(2002000), Emu(936000),
        )
    # Правый нижний — картинка для первого/последнего слайда
    if os.path.isfile(IMG_RIGHT_FIRST_LAST):
        slide.shapes.add_picture(
            IMG_RIGHT_FIRST_LAST,
            Emu(6876000), Emu(4060000),
            Emu(1858000), Emu(1725000),
        )

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 2 — О продукте (проблема + описание)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  О продукте", PINK, [
        ("1. Какую проблему решает продукт", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("Малый и средний бизнес в России активно использует Telegram для продаж "
         "и коммуникации с клиентами. По оценкам, более 50 000–70 000 предпринимателей "
         "ведут продажи через Telegram.", False, Pt(16), None),
        ("", False, Pt(4), None),
        ("Однако большинство из них не имеет технических навыков для создания "
         "автоматизированных ботов и вынуждены обрабатывать заявки вручную. "
         "Существующие платформы-конструкторы (BotHelp, SendPulse, Salebot) стоят "
         "от 10 000 до 70 000 ₽ и сложны в освоении.", False, Pt(16), None),
        ("", False, Pt(10), None),
        ("2. Описание конечного продукта", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("BotWork — сервис создания Telegram-ботов «под ключ»:", False, Pt(16), None),
        ("• Автоматическая воронка продаж (приветствие → презентация → оплата)",
         False, Pt(15), None),
        ("• Интеграция платежей (ЮKassa, СБП, Crypto)", False, Pt(15), None),
        ("• Настройка под нишу (курсы, товары, услуги)", False, Pt(15), None),
        ("• Визуальный конструктор для проектирования логики бота", False, Pt(15), None),
        ("• Стоимость: 5 000–25 000 ₽. Срок: 3–7 дней. Поддержка 2 нед. бесплатно.",
         False, Pt(15), RED),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 3 — О продукте (область применения)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  О продукте", PINK, [
        ("3. Область применения продукта", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("Способ применения:", True, Pt(16), GRAY),
        ("Бизнес получает готового Telegram-бота, который работает 24/7 внутри "
         "мессенджера — без отдельных приложений и сайтов.", False, Pt(16), None),
        ("", False, Pt(6), None),
        ("Ниши применения:", True, Pt(16), GRAY),
        ("• Образование — продажа курсов, вебинаров, консультаций через воронки",
         False, Pt(15), None),
        ("• Услуги — запись клиентов (салоны, репетиторы, консультанты), приём заявок",
         False, Pt(15), None),
        ("• Товары — мини-магазины с каталогом, корзиной и оплатой прямо в Telegram",
         False, Pt(15), None),
        ("", False, Pt(6), None),
        ("Выгоды для пользователей:", True, Pt(16), GRAY),
        ("• Автоматизация продаж и снижение нагрузки на персонал", False, Pt(15), None),
        ("• Увеличение конверсии на 15–25%", False, Pt(15), None),
        ("• Экономия на менеджерах", False, Pt(15), None),
        ("• Простота для клиентов — всё внутри привычного Telegram", False, Pt(15), None),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 4 — Технологичность (техническое решение)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Технологичность стартап-проекта", ORANGE, [
        ("1. Техническое решение проекта, инновационность", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("Технологический стек:", True, Pt(16), GRAY),
        ("• Backend: Python + фреймворк aiogram для асинхронной работы с Telegram Bot API",
         False, Pt(15), None),
        ("• Frontend-конструктор: веб-интерфейс (HTML/CSS/JS) для визуального "
         "проектирования логики бота методом drag & drop",
         False, Pt(15), None),
        ("• Модульная архитектура: набор готовых блоков — приветствие, кнопки, "
         "оплата, сбор данных, задержка, условия",
         False, Pt(15), None),
        ("• Интеграции: ЮKassa API, СБП, Crypto-платежи", False, Pt(15), None),
        ("", False, Pt(6), None),
        ("Инновационность:", True, Pt(16), GRAY),
        ("Уникальная гибридная модель: клиент самостоятельно проектирует логику бота "
         "в визуальном конструкторе, а затем передаёт схему разработчику для "
         "профессиональной реализации. Это сочетание no-code подхода с персонализированной "
         "разработкой «под ключ» — модель, которую не предлагают конкуренты.",
         False, Pt(15), None),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 5 — Технологичность (преимущества решения)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Технологичность стартап-проекта", ORANGE, [
        ("2. Преимущества выбранного технического решения", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("• Конструктор снижает порог входа: клиент описывает бота визуально, "
         "без технических знаний — это ускоряет коммуникацию и согласование",
         False, Pt(15), None),
        ("", False, Pt(3), None),
        ("• Шаблонный подход: типовые боты (визитка, заявки, магазин, FAQ) создаются "
         "за 2–3 дня, что обеспечивает высокую скорость вывода продукта",
         False, Pt(15), None),
        ("", False, Pt(3), None),
        ("• Низкая себестоимость: разработка одного бота занимает 2–3 дня при марже "
         "50–70%, что делает бизнес-модель устойчивой",
         False, Pt(15), None),
        ("", False, Pt(3), None),
        ("• Масштабируемость: модульная архитектура позволяет быстро адаптировать "
         "решение под новые ниши без переписывания кода",
         False, Pt(15), None),
        ("", False, Pt(3), None),
        ("• Персонализация: в отличие от массовых конструкторов, каждый бот "
         "настраивается индивидуально под конкретный бизнес клиента",
         False, Pt(15), None),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 6 — Технологичность (имеющийся задел)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Технологичность стартап-проекта", ORANGE, [
        ("3. Имеющийся задел для реализации проекта", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("Разработка (октябрь–ноябрь 2025):", True, Pt(16), GRAY),
        ("• Создано 2 MVP-бота для реальных клиентов (для коучей и товаров)",
         False, Pt(15), None),
        ("• Настроена интеграция с Telegram Bot API", False, Pt(15), None),
        ("• Запущена тестовая техподдержка", False, Pt(15), None),
        ("• Разработан визуальный конструктор чат-ботов (веб-интерфейс)",
         False, Pt(15), None),
        ("• Создан сайт-лендинг BotWork с описанием тарифов и конструктором",
         False, Pt(15), None),
        ("", False, Pt(6), None),
        ("Тестирование рынка:", True, Pt(16), GRAY),
        ("• Заведено 3 целевых Telegram-канала для продвижения", False, Pt(15), None),
        ("• Размещено 5 рекламных постов (бюджет 8 000 ₽)", False, Pt(15), None),
        ("• Получено 15 лидов, закрыто 2 сделки (конверсия 13%)", False, Pt(15), None),
        ("• Проведён анализ 5 конкурентов", False, Pt(15), None),
        ("• Сформирован финансовый план и юнит-экономика", False, Pt(15), None),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 7 — Технологичность (нацпроект + ИС)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Технологичность стартап-проекта", ORANGE, [
        ("4. Соответствие направлению национального проекта", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("Проект соответствует направлению «Экономика данных» национального проекта "
         "технологического лидерства:", False, Pt(16), None),
        ("• Развитие цифровых сервисов для бизнеса", False, Pt(15), None),
        ("• Автоматизация бизнес-процессов с помощью ИТ-технологий", False, Pt(15), None),
        ("• Повышение цифровой зрелости малого и среднего бизнеса", False, Pt(15), None),
        ("", False, Pt(10), None),
        ("5. Интеллектуальная собственность", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("На текущем этапе зарегистрированная ИС отсутствует.", False, Pt(16), None),
        ("Планируется регистрация программы для ЭВМ (визуальный конструктор ботов) "
         "в Роспатенте в течение первого года реализации проекта.",
         False, Pt(16), None),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 8 — Аналоги (существующие)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Аналоги и конкуренты", PURPLE, [
        ("1. Существующие аналоги", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("• BotHelp — конструктор чат-ботов и автоворонок. Средний чек 15 000–50 000 ₽. "
         "Требует самостоятельной настройки, сложный интерфейс.",
         False, Pt(15), None),
        ("", False, Pt(3), None),
        ("• SendPulse — маркетинговая платформа с модулем ботов. От 12 000 ₽/год. "
         "Широкий функционал, но избыточный для малого бизнеса.",
         False, Pt(15), None),
        ("", False, Pt(3), None),
        ("• Salebot — конструктор автоворонок в Telegram. От 10 000 ₽. "
         "Ориентирован на продвинутых пользователей.",
         False, Pt(15), None),
        ("", False, Pt(3), None),
        ("• PuzzleBot — визуальный конструктор ботов. От 990 ₽/мес. "
         "Ограниченный функционал в базовых тарифах.",
         False, Pt(15), None),
        ("", False, Pt(6), None),
        ("Все конкуренты предлагают инструменты для самостоятельной сборки ботов, "
         "но требуют значительных технических навыков или высокой абонентской платы. "
         "BotWork решает эту проблему, предлагая готовое решение «под ключ» по доступной цене.",
         False, Pt(15), RED),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 9 — Аналоги (конкурентные преимущества) — ТАБЛИЦА
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, "  Аналоги и конкуренты", PURPLE)

    add_textbox(slide, Emu(460000), Emu(750000), Emu(8200000), Emu(400000),
                [("2. Конкурентные преимущества продукта", True, Pt(20), DARK)])

    # Таблица
    rows, cols = 6, 5
    tbl_left = Emu(460000)
    tbl_top = Emu(1250000)
    tbl_w = Emu(8200000)
    tbl_h = Emu(3200000)
    table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h)
    table = table_shape.table

    headers = ["Критерий", "BotWork", "BotHelp", "SendPulse", "Salebot"]
    data = [
        ["Цена", "5–25 тыс. ₽", "15–50 тыс. ₽", "12 тыс.+ ₽/год", "10 тыс.+ ₽"],
        ["Персонализация", "Да, под ключ", "Нет, шаблоны", "Нет", "Частично"],
        ["Поддержка", "2 нед. бесплатно", "Платная", "Тикеты", "Платная"],
        ["Срок запуска", "3–7 дней", "Самостоятельно", "Самостоятельно", "Самостоятельно"],
        ["Конструктор", "Да", "Да", "Да", "Да"],
    ]

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_MAIN
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE

    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_MAIN
                p.font.size = Pt(11)
                p.font.bold = (j == 1)
                p.font.color.rgb = DARK

    # УЦП
    add_textbox(slide, Emu(460000), Emu(4600000), Emu(8200000), Emu(1800000),
                [
                    ("Уникальное ценностное предложение:", True, Pt(16), PURPLE),
                    ("Персонализированные Telegram-боты «под ключ» по цене в 2–3 раза ниже "
                     "рынка, с бесплатной поддержкой и запуском за 3–7 дней. Клиент получает "
                     "готовое решение, а не инструмент, который нужно осваивать.",
                     False, Pt(14), None),
                ])

    add_corner_images(slide)

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 10 — Востребованность
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Аналоги и конкуренты", PURPLE, [
        ("3. Востребованность продукта", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("Подтверждённый спрос:", True, Pt(16), GRAY),
        ("• 15 лидов получено через 3 Telegram-канала за первые 2 месяца работы",
         False, Pt(15), None),
        ("• 2 сделки закрыты (конверсия лид → продажа = 13%)", False, Pt(15), None),
        ("• Доход за стартовый период: 30 000 ₽", False, Pt(15), None),
        ("", False, Pt(6), None),
        ("Обратная связь от клиентов:", True, Pt(16), GRAY),
        ("• «Наконец-то доступная цена за бота — не нужно платить 50 тысяч»",
         False, Pt(15), None),
        ("• «Не нужно разбираться самому — просто объяснил что нужно и получил "
         "готовый бот»", False, Pt(15), None),
        ("", False, Pt(6), None),
        ("Объективные факторы:", True, Pt(16), GRAY),
        ("• 50 000–70 000 предпринимателей в РФ активно используют Telegram для продаж",
         False, Pt(15), None),
        ("• Рост интереса к автоматизации продаж в Telegram: +40% поисковых запросов "
         "за 2024–2025", False, Pt(15), None),
        ("• MAU Telegram в РФ превышает 100 млн пользователей", False, Pt(15), None),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 11 — Рынок (область применения)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Рынок и потребитель", BLUE, [
        ("1. Область применения продукта", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("Основные направления использования:", True, Pt(16), GRAY),
        ("", False, Pt(3), None),
        ("• Воронки продаж — автоматическая последовательность сообщений, "
         "ведущая клиента к покупке", False, Pt(15), None),
        ("• Запись на услуги — бот для салонов красоты, репетиторов, консультантов "
         "с выбором даты/времени", False, Pt(15), None),
        ("• Мини-магазины — каталог товаров с корзиной и онлайн-оплатой",
         False, Pt(15), None),
        ("• FAQ-боты — автоматические ответы на частые вопросы, снижение нагрузки "
         "на поддержку", False, Pt(15), None),
        ("• Рассылки — информирование клиентской базы о новинках и акциях",
         False, Pt(15), None),
        ("• Сбор заявок — автоматический приём и маршрутизация заявок менеджерам",
         False, Pt(15), None),
        ("• Чат-поддержка — первичная обработка обращений с передачей живому "
         "оператору при необходимости", False, Pt(15), None),
        ("", False, Pt(6), None),
        ("Дополнительные применения: геймификация, программы лояльности, "
         "агрегация отзывов, внутренние корпоративные боты.",
         False, Pt(14), GRAY),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 12 — Рынок (TAM-SAM-SOM)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Рынок и потребитель", BLUE, [
        ("2. Рынок, сегмент рынка, потенциальный потребитель", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("TAM (Total Addressable Market):", True, Pt(15), GRAY),
        ("~4,5 млн субъектов МСП в России. Потенциальный интерес к автоматизации "
         "через Telegram — ~500 000 компаний.", False, Pt(14), None),
        ("", False, Pt(3), None),
        ("SAM (Serviceable Available Market):", True, Pt(15), GRAY),
        ("~50 000–70 000 предпринимателей, активно использующих Telegram для продаж "
         "(Россия, СНГ). Ниши: образование, услуги, товары.", False, Pt(14), None),
        ("", False, Pt(3), None),
        ("SOM (Serviceable Obtainable Market):", True, Pt(15), GRAY),
        ("50–100 клиентов в первый год, 300+ клиентов к третьему году.",
         False, Pt(14), None),
        ("", False, Pt(6), None),
        ("Портрет потребителя:", True, Pt(16), GRAY),
        ("Владелец малого бизнеса (ИП/ООО), 25–45 лет, ведёт продажи через Telegram, "
         "оборот до 5 млн ₽/год. Ищет доступную автоматизацию без необходимости "
         "технических знаний.", False, Pt(14), None),
        ("", False, Pt(3), None),
        ("Тренды:", True, Pt(16), GRAY),
        ("Рост Telegram-коммерции +40% в год. MAU Telegram в РФ > 100 млн. "
         "Тренд на автоматизацию рутинных бизнес-процессов.", False, Pt(14), None),
        ("", False, Pt(3), None),
        ("География: Россия, СНГ (русскоязычная аудитория).", False, Pt(14), GRAY),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 13 — Рынок (барьеры)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Рынок и потребитель", BLUE, [
        ("3. Возможные барьеры для выхода на рынок", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("1. Высокая конкуренция", True, Pt(16), DARK),
        ("На рынке присутствуют крупные игроки (BotHelp, SendPulse) с большими бюджетами.",
         False, Pt(14), None),
        ("→ Преодоление: конкурируем ценой (в 2–3 раза ниже), персонализацией "
         "и качеством поддержки.", False, Pt(14), RED),
        ("", False, Pt(4), None),
        ("2. Низкий порог входа на рынок", True, Pt(16), DARK),
        ("Легко появляются новые мелкие конкуренты.", False, Pt(14), None),
        ("→ Преодоление: накопление клиентской базы, репутация, повторные продажи "
         "(доля повторных клиентов — до 30%).", False, Pt(14), RED),
        ("", False, Pt(4), None),
        ("3. Зависимость от Telegram API", True, Pt(16), DARK),
        ("Изменения в API могут потребовать доработки ботов.", False, Pt(14), None),
        ("→ Преодоление: мониторинг изменений, модульная архитектура для быстрой "
         "адаптации.", False, Pt(14), RED),
        ("", False, Pt(4), None),
        ("4. Ограниченный бюджет на маркетинг", True, Pt(16), DARK),
        ("→ Преодоление: бесплатные каналы (Telegram-каналы, холодные лиды), "
         "сотрудничество с микро-инфлюенсерами (2 000–5 000 ₽ за пост).",
         False, Pt(14), RED),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 14 — План реализации
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, "  План реализации проекта", RED)

    # Левая колонка — направления
    add_textbox(slide, Emu(460000), Emu(750000), Emu(4000000), Emu(5500000),
                [
                    ("Направления плана:", True, Pt(18), DARK),
                    ("", False, Pt(4), None),
                    ("1. Ключевая техническая задача", True, Pt(14), GRAY),
                    ("MVP → рассылки и аналитика → SaaS-конструктор", False, Pt(13), None),
                    ("", False, Pt(3), None),
                    ("2. Ресурсы проекта", True, Pt(14), GRAY),
                    ("Команда 3 чел., хостинг, домен, рекламный бюджет",
                     False, Pt(13), None),
                    ("", False, Pt(3), None),
                    ("3. Затраты на реализацию", True, Pt(14), GRAY),
                    ("0 мес: 20 000 ₽ | 3 мес: 45 000 ₽ | 12 мес: 60 000 ₽",
                     False, Pt(13), None),
                    ("", False, Pt(3), None),
                    ("4. Формирование команды", True, Pt(14), GRAY),
                    ("0 мес: 3 чел. | 12 мес: +1 разработчик", False, Pt(13), None),
                    ("", False, Pt(3), None),
                    ("5. Способ получения дохода", True, Pt(14), GRAY),
                    ("Разовая оплата 5–25 тыс. ₽ + доработки",
                     False, Pt(13), None),
                ])

    # Правая колонка — таймлайн
    add_textbox(slide, Emu(4700000), Emu(750000), Emu(4200000), Emu(5500000),
                [
                    ("Таймлайн:", True, Pt(18), DARK),
                    ("", False, Pt(4), None),
                    ("Старт (Q4 2025):", True, Pt(15), ORANGE),
                    ("2 MVP-бота, конструктор, сайт, 2 продажи",
                     False, Pt(13), None),
                    ("Доход: 30 000 ₽ | Расходы: 20 000 ₽", False, Pt(13), RED),
                    ("", False, Pt(4), None),
                    ("3 месяца (Q1 2026):", True, Pt(15), ORANGE),
                    ("Рассылки, аналитика, интеграция платежей, 8 ботов",
                     False, Pt(13), None),
                    ("Доход: 120 000 ₽ | Расходы: 45 000 ₽",
                     False, Pt(13), RED),
                    ("Конверсия: 15%", False, Pt(13), None),
                    ("", False, Pt(4), None),
                    ("12 месяцев (Q3 2026):", True, Pt(15), ORANGE),
                    ("Полная автоматизация, SaaS-конструктор, 25 ботов",
                     False, Pt(13), None),
                    ("Доход: 300 000 ₽ | Расходы: 60 000 ₽",
                     False, Pt(13), RED),
                    ("Конверсия: 25% | Повторные клиенты: 30%",
                     False, Pt(13), None),
                ])

    add_corner_images(slide)

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 15 — Квалификация (1/2)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Квалификация заявителя", RED, [
        ("1. Команда проекта, квалификация и опыт", True, Pt(20), DARK),
        ("", False, Pt(6), None),
        ("Артемий — Технический лидер", True, Pt(16), PURPLE),
        ("Разработка ботов, интеграции с API, техподдержка.",
         False, Pt(14), None),
        ("Дисциплины: программирование (Python), веб-разработка, базы данных.",
         False, Pt(14), GRAY),
        ("", False, Pt(4), None),
        ("Денис — Маркетинг и продажи", True, Pt(16), PURPLE),
        ("Продвижение, лидогенерация, переговоры с клиентами.",
         False, Pt(14), None),
        ("Дисциплины: маркетинг, управление проектами.",
         False, Pt(14), GRAY),
        ("", False, Pt(4), None),
        ("Борис — Старший разработчик", True, Pt(16), PURPLE),
        ("Разработка конструктора чат-ботов, ботов.",
         False, Pt(14), None),
        ("Дисциплины: программирование, UX/UI-дизайн.",
         False, Pt(14), GRAY),
        ("", False, Pt(6), None),
        ("2. Опыт реализации проектов по теме заявки", True, Pt(20), DARK),
        ("", False, Pt(4), None),
        ("• Создано 2 MVP-бота для реальных клиентов", False, Pt(14), None),
        ("• Разработан визуальный конструктор чат-ботов", False, Pt(14), None),
        ("• Закрыто 2 сделки, доход 30 000 ₽", False, Pt(14), None),
        ("• Проведён анализ рынка и 5 конкурентов", False, Pt(14), None),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 16 — Квалификация (2/2)
    # ══════════════════════════════════════════════════════════════════
    content_slide(prs, "  Квалификация заявителя", RED, [
        ("3. Участие в программе «Стартап как диплом»", True, Pt(20), DARK),
        ("", False, Pt(4), None),
        ("Участие в образовательных программах в области предпринимательства.",
         False, Pt(15), None),
        ("[Указать майнор, если применимо]", False, Pt(14), GRAY),
        ("", False, Pt(8), None),
        ("4. Участие в акселераторах, бизнес-инкубаторах", True, Pt(20), DARK),
        ("", False, Pt(4), None),
        ("[Указать акселератор HSE, если применимо]", False, Pt(14), GRAY),
        ("", False, Pt(8), None),
        ("5. Участие в ПИШ / молодежной лаборатории / СКБ", True, Pt(20), DARK),
        ("", False, Pt(4), None),
        ("[Указать, если применимо]", False, Pt(14), GRAY),
    ])

    # ══════════════════════════════════════════════════════════════════
    # СЛАЙД 17 — Спасибо за внимание
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Декоративные блоки (зеркально к титульному)
    add_rect(slide, Emu(3000000), Emu(0), Emu(6200000), Emu(1800000), BLUE)
    add_rect(slide, Emu(5300000), Emu(0), Emu(3900000), Emu(1800000), PURPLE)
    add_rect(slide, Emu(3000000), Emu(1800000), Emu(6200000), Emu(5100000), PINK)
    add_rect(slide, Emu(5300000), Emu(1800000), Emu(3900000), Emu(5100000), ORANGE)
    add_rect(slide, Emu(3000000), Emu(5300000), Emu(2300000), Emu(1560000), RED)

    # «Спасибо за внимание!»
    add_textbox(slide, Emu(5500000), Emu(200000), Emu(3500000), Emu(1200000),
                [
                    ("Спасибо", True, Pt(36), WHITE),
                    ("за внимание!", True, Pt(36), WHITE),
                ],
                alignment=PP_ALIGN.CENTER)

    # Контактные данные
    add_textbox(slide, Emu(800000), Emu(1900000), Emu(2000000), Emu(800000),
                [
                    ("BotWork", True, Pt(28), DARK),
                    ("Номер заявки: _______", False, Pt(14), DARK),
                    ("ФИО: Репин Артемий, Мартынов Борис, Харисов Денис", False, Pt(14), DARK),
                ])

    add_textbox(slide, Emu(5500000), Emu(2000000), Emu(3500000), Emu(1600000),
                [
                    ("Почта", True, Pt(16), DARK),
                    ("[email]", False, Pt(14), GRAY),
                    ("", False, Pt(6), None),
                    ("Телефон", True, Pt(16), DARK),
                    ("[телефон]", False, Pt(14), GRAY),
                    ("", False, Pt(6), None),
                    ("Telegram", True, Pt(16), DARK),
                    ("@borysikkk", False, Pt(14), GRAY),
                ])

    add_corner_images(slide, first_last=True)

    # ─── Сохранение ───────────────────────────────────────────────────
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUTPUT_DIR / "BotWork_Application.pptx"
    prs.save(str(out))
    print(f"Презентация сохранена: {out}")


if __name__ == "__main__":
    create_presentation()
